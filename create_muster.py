from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

def create_muster_template():
    prs = Presentation()
    
    # Placeholder colors
    bg_color = RGBColor(240, 240, 240)
    primary_color = RGBColor(0, 102, 204)
    text_color = RGBColor(51, 51, 51)
    
    # Title and 9 Content slides
    slide_layouts = [0, 1, 1, 1, 1, 1, 1, 1, 1, 1]
    
    for i, layout_idx in enumerate(slide_layouts):
        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        if slide.shapes.title:
            slide.shapes.title.text = f"Slide {i+1} Title"
            for paragraph in slide.shapes.title.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "Arial"
                    run.font.color.rgb = primary_color
                    
        # Content
        if len(slide.shapes.placeholders) > 1:
            try:
                body_shape = slide.shapes.placeholders[1]
                tf = body_shape.text_frame
                tf.text = f"This is placeholder content for Slide {i+1}."
                p = tf.add_paragraph()
                p.text = "More placeholder text here."
                
                for paragraph in tf.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = "Arial"
                        run.font.color.rgb = text_color
            except Exception:
                pass
                
    prs.save("Muster_Template.pptx")
    print("Muster_Template.pptx created successfully.")

if __name__ == "__main__":
    create_muster_template()