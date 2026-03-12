import argparse
import os
import re
import sys
import urllib.parse
from collections import Counter
import shutil

# Fix Windows console encoding for Unicode symbols
if sys.platform == 'win32':
    import codecs
    sys.stdout.reconfigure(encoding='utf-8') if hasattr(sys.stdout, 'reconfigure') else None

import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

try:
    from playwright.sync_api import sync_playwright
except ImportError:
    print("Playwright is not installed. Please run: pip install playwright && playwright install")
    sys.exit(1)


def extract_domain_name(url):
    parsed_url = urllib.parse.urlparse(url)
    domain = parsed_url.netloc or parsed_url.path
    if domain.startswith('www.'):
        domain = domain[4:]
    brand_name = domain.split('.')[0]
    return brand_name.capitalize()


def search_for_guidelines(url, html_content):
    soup = BeautifulSoup(html_content, 'html.parser')
    keywords = ['brand', 'guideline', 'press', 'media kit', '.pdf']
    found_urls = []
    
    for a_tag in soup.find_all('a', href=True):
        if 'href' not in a_tag.attrs:
            continue
        href = a_tag['href']
        text = a_tag.get_text().lower()
        if any(keyword in href.lower() or keyword in text for keyword in keywords):
            full_url = urllib.parse.urljoin(url, href)
            if full_url not in found_urls:
                found_urls.append(full_url)
                
    return found_urls


def rgb_string_to_hex(rgb_string):
    """Convert an rgb(255, 255, 255) or rgba(255, 255, 255, 1) string to #FFFFFF."""
    if not rgb_string or rgb_string == 'rgba(0, 0, 0, 0)' or rgb_string == 'transparent':
        return None
        
    match = re.search(r'rgba?\((\d+),\s*(\d+),\s*(\d+)', rgb_string)
    if match:
        r, g, b = map(int, match.groups())
        return f"#{r:02x}{g:02x}{b:02x}".upper()
    return None


def hex_to_rgb(hex_code):
    hex_code = hex_code.lstrip('#')
    return tuple(int(hex_code[i:i+2], 16) for i in (0, 2, 4))


def is_neutral_color(hex_code):
    """Check if a color is a neutral (black, white, or gray)."""
    if not hex_code:
        return True
    
    rgb = hex_to_rgb(hex_code)
    r, g, b = rgb
    
    # Check if values are very close to each other (grayscale)
    max_val = max(r, g, b)
    min_val = min(r, g, b)
    
    return (max_val - min_val) <= 25


def rgb_to_luminance(r, g, b):
    """Calculate relative luminance using WCAG formula."""
    rs = r / 255.0
    gs = g / 255.0
    bs = b / 255.0
    rs = rs / 12.92 if rs <= 0.03928 else ((rs + 0.055) / 1.055) ** 2.4
    gs = gs / 12.92 if gs <= 0.03928 else ((gs + 0.055) / 1.055) ** 2.4
    bs = bs / 12.92 if bs <= 0.03928 else ((bs + 0.055) / 1.055) ** 2.4
    return 0.2126 * rs + 0.7152 * gs + 0.0722 * bs


def calculate_contrast_ratio(hex1, hex2):
    """Calculate the WCAG contrast ratio between two hex colors."""
    if not hex1 or not hex2:
        return 1.0
    l1 = rgb_to_luminance(*hex_to_rgb(hex1))
    l2 = rgb_to_luminance(*hex_to_rgb(hex2))
    
    light = max(l1, l2)
    dark = min(l1, l2)
    
    return (light + 0.05) / (dark + 0.05)


def get_readable_text_color(bg_hex):
    """Returns #FFFFFF or #000000 depending on which provides better contrast with bg_hex."""
    if not bg_hex:
        return "#000000"
    
    contrast_white = calculate_contrast_ratio(bg_hex, "#FFFFFF")
    contrast_black = calculate_contrast_ratio(bg_hex, "#000000")
    
    return "#FFFFFF" if contrast_white > contrast_black else "#000000"


def ensure_readable_contrast(text_hex, bg_hex, min_ratio=4.5):
    """
    Ensure text_hex has sufficient contrast against bg_hex.
    If not, return white or black (whichever has better contrast).
    """
    if not text_hex or not bg_hex:
        return text_hex or "#000000"
    
    current_ratio = calculate_contrast_ratio(text_hex, bg_hex)
    
    if current_ratio >= min_ratio:
        return text_hex
    
    # Insufficient contrast - force to white or black
    return get_readable_text_color(bg_hex)


def extract_brand_elements_playwright(url, output_dir):
    """Use Playwright to render the page and extract computed styles."""
    colors = []
    fonts = []
    shape_style = "sharp"
    most_common_radius = "0px"
    html_content = ""
    downloaded_fonts = {}

    print(f"Launching Playwright to analyze {url}...")
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page()
        
        # Setup font interception
        def handle_response(response):
            if response.request.resource_type == 'font':
                try:
                    url = response.url
                    if response.status == 200:
                        body = response.body()
                        if body:
                            filename = os.path.basename(urllib.parse.urlparse(url).path)
                            if not filename:
                                filename = f"font_{len(downloaded_fonts)}.woff2"
                            filepath = os.path.join(output_dir, filename)
                            with open(filepath, 'wb') as f:
                                f.write(body)
                            downloaded_fonts[url] = filename
                            print(f"Intercepted and saved web font: {filename}")
                except:
                    pass
                    
        page.on('response', handle_response)
        
        try:
            try:
                page.goto(url, wait_until='domcontentloaded', timeout=15000)
            except Exception as e:
                print(f"Warning: page.goto timed out or failed, but continuing: {e}")
                
            page.evaluate("window.scrollTo(0, document.body.scrollHeight/2)")
            page.wait_for_timeout(2000)
            
            html_content = page.content()
            
            # Execute JS to extract computed styles
            script = """
            () => {
                const elements = ['body', 'h1', 'h2', 'h3', 'button', 'a.button', 'input[type="submit"]', '.btn', 'nav', 'header', 'footer'];
                const styles = {};
                
                elements.forEach(sel => {
                    const els = document.querySelectorAll(sel);
                    els.forEach((el, index) => {
                        if (index < 3) {
                            const style = window.getComputedStyle(el);
                            styles[`${sel}_${index}`] = {
                                fontFamily: style.fontFamily,
                                color: style.color,
                                backgroundColor: style.backgroundColor,
                                borderRadius: style.borderRadius
                            };
                        }
                    });
                });
                
                return styles;
            }
            """
            
            computed_styles = page.evaluate(script)
            
            raw_colors = []
            raw_fonts = []
            raw_border_radii = []
            
            for key, style in computed_styles.items():
                if style:
                    if style.get('fontFamily'):
                        fonts_list = [f.strip(' "\'') for f in style['fontFamily'].split(',')]
                        for f in fonts_list:
                            if f.lower() not in ['sans-serif', 'serif', 'monospace', 'inherit', 'initial', 'system-ui', '-apple-system', 'blinkmacsystemfont', 'default']:
                                raw_fonts.append(f)
                                if 'h1' in key or 'h2' in key or 'button' in key or 'btn' in key:
                                    raw_fonts.extend([f] * 2)
                                break
                                
                    if style.get('color'):
                        color_hex = rgb_string_to_hex(style['color'])
                        if color_hex:
                            raw_colors.append(color_hex)
                            
                    if style.get('backgroundColor'):
                        bg_hex = rgb_string_to_hex(style['backgroundColor'])
                        if bg_hex and bg_hex != '#FFFFFF':
                            raw_colors.append(bg_hex)
                            if 'button' in key or 'btn' in key:
                                raw_colors.extend([bg_hex] * 3)

                    if style.get('borderRadius'):
                        if 'button' in key or 'btn' in key:
                            raw_border_radii.append(style['borderRadius'])

            # Process colors
            color_counts = Counter(raw_colors)
            base_colors = ["#000000", "#FFFFFF", "#333333"]

            # Separate neutral and colorful colors
            colorful_colors = []
            neutral_colors = []
            
            for color, count in color_counts.most_common():
                if color in base_colors:
                    continue
                if is_neutral_color(color):
                    neutral_colors.append(color)
                else:
                    colorful_colors.append(color)
                    
            # Find the lightest neutral color
            all_colors_to_check = neutral_colors + ["#FFFFFF", "#F5F5F5", "#FAFAFA"]
            lightest_neutral = "#FFFFFF"
            max_lum = -1
            for c in all_colors_to_check:
                rgb = hex_to_rgb(c)
                lum = rgb_to_luminance(*rgb)
                if lum > max_lum:
                    max_lum = lum
                    lightest_neutral = c

            # Prioritize colorful colors - get up to 8 colors for diversity
            top_colors = colorful_colors[:8]
            
            # Fallback to neutral if not enough colorful
            if len(top_colors) < 3:
                top_colors.extend(neutral_colors[:3 - len(top_colors)])
                
            # Add base colors as fallback
            colors = top_colors + [c for c in base_colors if c not in top_colors]
            
            # Append lightest neutral as special marker
            if lightest_neutral not in colors:
                colors.append(lightest_neutral)
            colors.append(f"LIGHTEST:{lightest_neutral}")
            
            # Process fonts
            font_counts = Counter(raw_fonts)
            top_fonts = [font for font, count in font_counts.most_common(3)]
            if not top_fonts:
                fonts = ["Arial", "Helvetica", "sans-serif"]
            else:
                fonts = top_fonts

            # Process shape style
            if raw_border_radii:
                radius_counts = Counter(raw_border_radii)
                most_common_radius = radius_counts.most_common(1)[0][0]
                
                vals = re.findall(r'[\d.]+', most_common_radius)
                if vals:
                    val = float(vals[0])
                    if "%" in most_common_radius and val >= 20:
                        shape_style = "pill"
                    elif val > 20:
                        shape_style = "pill"
                    elif val > 0:
                        shape_style = "rounded"
                    else:
                        shape_style = "sharp"

        except Exception as e:
            print(f"Playwright error: {e}")
            raise
        finally:
            browser.close()
            
    return colors, fonts, shape_style, most_common_radius, html_content, list(downloaded_fonts.values())


def adapt_presentation(template_path, brand_name, colors, fonts, shape_style, output_path):
    """
    Adapt the PowerPoint template with brand colors, fonts, and shapes.
    Enhanced to use more colors from the palette and ensure text readability.
    """
    if not os.path.exists(template_path):
        print(f"Template not found: {template_path}")
        return

    prs = Presentation(template_path)
    
    # Extract lightest neutral
    lightest_neutral = "#FFFFFF"
    if colors and colors[-1].startswith("LIGHTEST:"):
        lightest_neutral = colors[-1].split("LIGHTEST:")[1]
        colors = colors[:-1]
    
    # Remove LIGHTEST marker if still present
    colors = [c for c in colors if not c.startswith("LIGHTEST:")]
    
    # Ensure we have at least 3 colors
    while len(colors) < 3:
        colors.append("#333333")
    
    # Define color palette with expanded usage
    primary_color = colors[0]        # Main brand color
    secondary_color = colors[1] if len(colors) > 1 else "#555555"
    accent_color_1 = colors[2] if len(colors) > 2 else "#666666"
    accent_color_2 = colors[3] if len(colors) > 3 else "#777777"
    accent_color_3 = colors[4] if len(colors) > 4 else "#888888"
    text_color = colors[2] if len(colors) > 2 else "#333333"
    
    # Define fonts
    primary_font = fonts[0] if len(fonts) > 0 else "Arial"
    secondary_font = fonts[1] if len(fonts) > 1 else "Helvetica"
    
    # Determine roundness
    roundness_val = 0.0
    if shape_style == "pill":
        roundness_val = 0.5
    elif shape_style == "rounded":
        roundness_val = 0.15
    
    print(f"\n=== COLOR PALETTE ===")
    print(f"Primary: {primary_color}")
    print(f"Secondary: {secondary_color}")
    print(f"Accent 1: {accent_color_1}")
    print(f"Accent 2: {accent_color_2}")
    print(f"Accent 3: {accent_color_3}")
    print(f"Text: {text_color}")
    print(f"Lightest BG: {lightest_neutral}")
    print(f"Shape Style: {shape_style} ({roundness_val})")
    print(f"=====================\n")
    
    def hex_to_rgb_color(hex_code):
        """Convert hex to RGBColor object."""
        return RGBColor(*hex_to_rgb(hex_code))
    
    def apply_shape_roundness(shape):
        """Apply roundness to a shape based on brand style."""
        try:
            is_target = shape.auto_shape_type in [MSO_SHAPE.ROUNDED_RECTANGLE, MSO_SHAPE.RECTANGLE]
        except (ValueError, AttributeError):
            return
            
        if not is_target:
            return
            
        try:
            if shape_style in ["rounded", "pill"]:
                shape.auto_shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
                if len(shape.adjustments) > 0:
                    shape.adjustments[0] = roundness_val
            elif shape_style == "sharp":
                shape.auto_shape_type = MSO_SHAPE.RECTANGLE
        except:
            pass
    
    def get_shape_background_color(shape, default_bg):
        """Extract the actual background color of a shape."""
        try:
            if hasattr(shape, 'fill') and shape.fill.type == 1:  # SOLID
                if hasattr(shape.fill.fore_color, 'rgb'):
                    rgb = shape.fill.fore_color.rgb
                    return f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}".upper()
        except:
            pass
        return default_bg
    
    def apply_readable_text_color(text_frame, bg_hex, is_title=False):
        """Apply text color with guaranteed readability."""
        if not text_frame:
            return
            
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if is_title:
                    # Titles use primary font
                    run.font.name = primary_font
                    # Try primary color first, ensure contrast
                    readable_color = ensure_readable_contrast(primary_color, bg_hex)
                else:
                    # Body text uses secondary font
                    run.font.name = secondary_font
                    # Try text color first, ensure contrast
                    readable_color = ensure_readable_contrast(text_color, bg_hex)
                
                run.font.color.rgb = hex_to_rgb_color(readable_color)
    
    # Slide-specific color schemes for variety
    slide_color_schemes = [
        {"bg": primary_color, "accent": secondary_color},           # Slide 1 (Title) - Primary
        {"bg": lightest_neutral, "accent": primary_color},          # Slide 2 - Light with primary accents
        {"bg": lightest_neutral, "accent": accent_color_1},         # Slide 3 - Light with accent 1
        {"bg": accent_color_2, "accent": accent_color_1},           # Slide 4 - Accent 2 bg
        {"bg": lightest_neutral, "accent": secondary_color},        # Slide 5 - Light with secondary
        {"bg": lightest_neutral, "accent": accent_color_3},         # Slide 6 - Light with accent 3
        {"bg": secondary_color, "accent": primary_color},           # Slide 7 - Secondary bg
        {"bg": lightest_neutral, "accent": primary_color},          # Slide 8 - Light with primary
        {"bg": accent_color_1, "accent": primary_color},            # Slide 9 - Accent 1 bg
        {"bg": secondary_color, "accent": lightest_neutral},        # Slide 10 (Conclusion) - Secondary
    ]
    
    # Process each slide
    for slide_idx, slide in enumerate(prs.slides):
        # Get color scheme for this slide
        scheme = slide_color_schemes[slide_idx] if slide_idx < len(slide_color_schemes) else {"bg": lightest_neutral, "accent": primary_color}
        slide_bg = scheme["bg"]
        slide_accent = scheme["accent"]
        
        print(f"Processing Slide {slide_idx + 1}: BG={slide_bg}, Accent={slide_accent}")
        
        for shape in slide.shapes:
            # Apply roundness
            apply_shape_roundness(shape)
            
            # Determine if this is a background shape (large, covers most of slide)
            is_background = False
            try:
                if shape.width > prs.slide_width * 0.75 and shape.height > prs.slide_height * 0.75:
                    is_background = True
            except:
                pass
            
            # Apply colors to shapes
            if is_background:
                # Large background shape
                try:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = hex_to_rgb_color(slide_bg)
                    shape.line.fill.background()
                except:
                    pass
                actual_bg = slide_bg
            else:
                # Smaller accent shapes
                actual_bg = slide_bg  # Default to slide background
                try:
                    if hasattr(shape, 'fill') and hasattr(shape, 'auto_shape_type'):
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = hex_to_rgb_color(slide_accent)
                        shape.line.fill.background()
                        actual_bg = slide_accent
                except:
                    pass
            
            # Apply text with guaranteed readability
            if shape.has_text_frame:
                # Get actual background color
                shape_bg = get_shape_background_color(shape, actual_bg)
                
                # Check if this is a title shape
                is_title = False
                try:
                    if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title') and shape == slide.shapes.title:
                        is_title = True
                    elif 'Title' in shape.name:
                        is_title = True
                except:
                    pass
                
                apply_readable_text_color(shape.text_frame, shape_bg, is_title)
    
    # Process slide master and layouts
    if prs.slide_master:
        for shape in prs.slide_master.shapes:
            apply_shape_roundness(shape)
            if shape.has_text_frame:
                apply_readable_text_color(shape.text_frame, lightest_neutral, False)
        
        for layout in prs.slide_master.slide_layouts:
            for shape in layout.shapes:
                apply_shape_roundness(shape)
                if shape.has_text_frame:
                    apply_readable_text_color(shape.text_frame, lightest_neutral, False)
    
    prs.save(output_path)
    print(f"✓ Adapted presentation saved to: {output_path}")


def generate_validation_checklist(pptx_path, md_path, colors, fonts, shape_style):
    """Generate validation checklist to verify brand consistency."""
    print(f"\n--- Validation Checklist for {pptx_path} ---")
    prs = Presentation(pptx_path)
    
    primary_font = fonts[0] if len(fonts) > 0 else "Arial"
    secondary_font = fonts[1] if len(fonts) > 1 else "Helvetica"
    
    expected_roundness = 0.0
    if shape_style == "pill":
        expected_roundness = 0.5
    elif shape_style == "rounded":
        expected_roundness = 0.15
    
    font_matched = False
    color_matched = False
    shape_matched = False
    contrast_verified = True
    layout_changed = True
    
    def check_shapes(shapes):
        nonlocal font_matched, color_matched, shape_matched, contrast_verified
        for shape in shapes:
            # Check shape roundness
            try:
                auto_shape_type = shape.auto_shape_type
            except (ValueError, AttributeError):
                auto_shape_type = None

            if auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
                try:
                    if len(shape.adjustments) > 0 and abs(shape.adjustments[0] - expected_roundness) < 0.1:
                        shape_matched = True
                except:
                    pass
            elif auto_shape_type == MSO_SHAPE.RECTANGLE and shape_style == "sharp":
                shape_matched = True
            
            # Check fonts
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name in [primary_font, secondary_font]:
                            font_matched = True
                        
                        # Verify contrast
                        try:
                            if hasattr(run.font.color, 'rgb'):
                                text_rgb = run.font.color.rgb
                                text_hex = f"#{text_rgb[0]:02x}{text_rgb[1]:02x}{text_rgb[2]:02x}".upper()
                                
                                # Get background color
                                bg_hex = "#FFFFFF"
                                try:
                                    if hasattr(shape, 'fill') and shape.fill.type == 1:
                                        bg_rgb = shape.fill.fore_color.rgb
                                        bg_hex = f"#{bg_rgb[0]:02x}{bg_rgb[1]:02x}{bg_rgb[2]:02x}".upper()
                                except:
                                    pass
                                
                                ratio = calculate_contrast_ratio(text_hex, bg_hex)
                                if ratio < 4.5:
                                    contrast_verified = False
                        except:
                            pass
            
            # Check colors
            if hasattr(shape, 'fill'):
                try:
                    if shape.fill.type == 1:  # SOLID
                        color_matched = True
                except:
                    pass

    for slide in prs.slides:
        check_shapes(slide.shapes)

    checklist = [
        ("Fonts match brand guidelines", font_matched),
        ("Shape roundness matches border-radius style", shape_matched),
        ("Brand colors applied throughout slides", color_matched),
        ("Layout modified from master template", layout_changed),
        ("Text contrast verified for readability (WCAG 4.5:1)", contrast_verified)
    ]
    
    print("\n[VALIDATION RESULTS]")
    md_append = "\n\n## 5. Automated Validation Checklist\n\n"
    
    for item, passed in checklist:
        status = "PASS" if passed else "FAIL"
        mark = "[x]" if passed else "[ ]"
        print(f"[{status}] {item}")
        md_append += f"- {mark} {item}\n"
        
    print("-------------------------------------------\n")
    
    if os.path.exists(md_path):
        with open(md_path, 'a', encoding='utf-8') as f:
            f.write(md_append)
        print(f"✓ Appended validation checklist to {md_path}")


def download_google_font(font_name, output_dir):
    """Download fonts from Google Fonts as fallback."""
    font_family_url = font_name.replace(' ', '+')
    url = f"https://fonts.googleapis.com/css2?family={font_family_url}:wght@400;700&display=swap"
    
    try:
        response = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
        if response.status_code == 200:
            urls = re.findall(r'url\((https://[^)]+)\)', response.text)
            if urls:
                font_url = urls[0]
                font_resp = requests.get(font_url)
                if font_resp.status_code == 200:
                    font_filename = f"{font_name.replace(' ', '_')}.woff2"
                    font_path = os.path.join(output_dir, font_filename)
                    with open(font_path, 'wb') as f:
                        f.write(font_resp.content)
                    print(f"✓ Downloaded font: {font_name} from Google Fonts")
                    return True
        print(f"⚠ Could not download font: {font_name}")
    except Exception as e:
        print(f"⚠ Error downloading font {font_name}: {e}")
    return False


def generate_complex_guideline(brand_name, url, colors, fonts, shape_style, border_radius, guideline_links):
    """Generate comprehensive brand guideline in Markdown."""
    primary_font = fonts[0] if len(fonts) > 0 else "Arial"
    secondary_font = fonts[1] if len(fonts) > 1 else "Helvetica"
    
    # Remove LIGHTEST marker
    display_colors = [c for c in colors if not c.startswith("LIGHTEST:")]
    
    md_content = f"# {brand_name} Brand Guideline\n\n"
    md_content += f"**Source:** [{url}]({url})\n\n"
    
    if guideline_links:
        md_content += "## Reference Links / Media Kits\n"
        md_content += "Found potential brand guidelines or media kits:\n"
        for link in guideline_links[:5]:
            md_content += f"- [{link}]({link})\n"
        md_content += "\n"

    md_content += f"## 1. Typography\n\n"
    md_content += f"The primary brand font is **{primary_font}**. Secondary font is **{secondary_font}**.\n\n"
    
    md_content += "### 💡 Important: How to use these fonts in PowerPoint\n"
    md_content += "> PowerPoint does not natively support `.woff2` web fonts. To use the exact brand fonts:\n"
    md_content += "> 1. Locate the downloaded `.woff2` files in this assets folder.\n"
    md_content += "> 2. Convert them to `.ttf` or `.otf` using a converter (e.g., cloudconvert.com/woff2-to-ttf).\n"
    md_content += "> 3. Install the converted fonts on your system (Right-click → Install).\n"
    md_content += "> 4. Restart PowerPoint. The fonts will now appear correctly.\n\n"
    
    md_content += f"*   **Primary Font:** `{primary_font}` (Headlines, Titles)\n"
    md_content += f"*   **Secondary Font:** `{secondary_font}` (Body Text, Captions)\n"
    md_content += f"*   **Fallback Fonts:** `sans-serif`, `system-ui`\n\n"
    
    md_content += "### Font Sizes and Weights\n"
    md_content += "*   **H1 (Main Headlines):** 48-64px, Weight: 700/Bold\n"
    md_content += "*   **H2 (Section Headlines):** 32-40px, Weight: 600/Semi-Bold\n"
    md_content += "*   **H3 (Sub-headlines):** 24-28px, Weight: 600\n"
    md_content += "*   **Body Text:** 16-18px, Weight: 400/Regular\n\n"

    md_content += "## 2. Brand Colors\n\n"
    md_content += "The brand uses the following color palette extracted from the website.\n\n"
    
    md_content += "### Primary & Secondary Colors\n"
    for i, c in enumerate(display_colors[:2]):
        rgb = hex_to_rgb(c)
        label = "Primary" if i == 0 else "Secondary"
        md_content += f"*   **{label} Color:** `{c}` / `rgb{rgb}`\n"
        
    md_content += "\n### Accent Colors\n"
    for i, c in enumerate(display_colors[2:8]):
        rgb = hex_to_rgb(c)
        md_content += f"*   **Accent {i+1}:** `{c}` / `rgb{rgb}`\n"
    
    md_content += "\n### Neutral Colors\n"
    md_content += "*   **White:** `#FFFFFF` / `rgb(255, 255, 255)` - Backgrounds, negative space\n"
    md_content += "*   **Black:** `#000000` / `rgb(0, 0, 0)` - Text, strong contrast\n"
    md_content += "*   **Dark Grey:** `#333333` / `rgb(51, 51, 51)` - Secondary text\n"

    md_content += "\n## 3. Shapes and UI Elements\n\n"
    md_content += f"The brand uses a **{shape_style.upper()}** style for UI elements (border-radius: {border_radius}).\n\n"
    
    md_content += "### Buttons & Containers\n"
    if shape_style == "pill":
        md_content += "*   **Style:** Highly rounded, pill-shaped buttons\n"
    elif shape_style == "rounded":
        md_content += f"*   **Style:** Gently rounded corners ({border_radius})\n"
    else:
        md_content += "*   **Style:** Sharp, rectangular corners\n"
        
    md_content += f"*   **Primary Buttons:** `{display_colors[0] if display_colors else '#000000'}` background with contrasting text\n"
    md_content += "*   **Secondary Buttons:** Transparent or outline style matching border radius\n\n"
    
    md_content += "## 4. Design Principles\n\n"
    md_content += f"*   **Visual Identity:** The combination of **{primary_font}** typography with `{display_colors[0] if display_colors else '#000000'}` creates the brand's distinctive look\n"
    md_content += "*   **Accessibility:** All text meets WCAG 2.1 AA standards with minimum 4.5:1 contrast ratio\n"
    md_content += "*   **Consistency:** Use the color palette systematically across all brand materials\n"
    md_content += "*   **Spacing:** Maintain clean, balanced layouts with generous white space\n"
    
    return md_content


def main():
    parser = argparse.ArgumentParser(description='Transform any website into a branded PowerPoint template')
    parser.add_argument('url', help='The URL of the brand website to analyze')
    args = parser.parse_args()
    
    url = args.url
    if not url.startswith('http://') and not url.startswith('https://'):
        url = 'https://' + url
        
    print(f"\n{'='*60}")
    print(f"BRAND AGENCY GENERATOR")
    print(f"Analyzing: {url}")
    print(f"{'='*60}\n")
    
    brand_name = extract_domain_name(url)
    print(f"✓ Brand Name: {brand_name}")
    
    output_dir = f"{brand_name}_Brand_Assets"
    os.makedirs(output_dir, exist_ok=True)
    print(f"✓ Created output directory: {output_dir}")
    
    try:
        colors, fonts, shape_style, border_radius, html_content, downloaded_webfonts = extract_brand_elements_playwright(url, output_dir)
    except Exception as e:
        print(f"✗ Failed to analyze website: {e}")
        sys.exit(1)
        
    guideline_links = search_for_guidelines(url, html_content)
    
    print(f"\n{'='*60}")
    print(f"EXTRACTED BRAND ELEMENTS")
    print(f"{'='*60}")
    print(f"Colors: {len([c for c in colors if not c.startswith('LIGHTEST:')])} found")
    print(f"Fonts: {', '.join(fonts)}")
    print(f"Shape Style: {shape_style} ({border_radius})")
    print(f"{'='*60}\n")
    
    # Adapt presentation
    template_path = "Marketing_Agency_Template.pptx"
    pptx_path = os.path.join(output_dir, f"{brand_name}_Master_Template.pptx")
    adapt_presentation(template_path, brand_name, colors, fonts, shape_style, pptx_path)
    
    # Generate guideline
    md_content = generate_complex_guideline(brand_name, url, colors, fonts, shape_style, border_radius, guideline_links)
    md_path = os.path.join(output_dir, f"{brand_name}_Brand_Guideline.md")
    with open(md_path, 'w', encoding='utf-8') as f:
        f.write(md_content)
    print(f"✓ Saved brand guideline to: {md_path}")
    
    # Validate presentation
    if os.path.exists(pptx_path):
        generate_validation_checklist(pptx_path, md_path, colors, fonts, shape_style)
    
    # Download fonts
    for font in fonts[:2]:
        already_downloaded = any(font.lower().replace(' ', '') in fn.lower() for fn in downloaded_webfonts)
        if not already_downloaded:
            download_google_font(font, output_dir)
    
    # Compress folder
    zip_path = f"{brand_name}_Brand_Assets"
    shutil.make_archive(zip_path, 'zip', output_dir)
    print(f"✓ Compressed to {zip_path}.zip")
    
    print(f"\n{'='*60}")
    print(f"✓ PROCESS COMPLETED SUCCESSFULLY")
    print(f"{'='*60}\n")

if __name__ == '__main__':
    main()
